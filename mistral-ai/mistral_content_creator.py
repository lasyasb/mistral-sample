"""
Mistral Content Creator - Blog/Social/Email + PPTX Slide Generator
"""
import os
import json
import time
import requests
import argparse
from datetime import datetime
from dotenv import load_dotenv
from pptx import Presentation
from pptx.util import Inches

# Load API key from .env
load_dotenv()
API_KEY = os.getenv("MISTRAL_API_KEY")
print("🔑 Loaded API Key:", API_KEY)

HEADERS = {"Authorization": f"Bearer {API_KEY}", "Content-Type": "application/json"}

def generate_content_streaming(topic, content_type="blog", tone="professional"):
    """Generate content using Mistral AI with streaming output."""
    endpoint = "https://api.mistral.ai/v1/chat/completions"

    # Prompts
    prompts = {
        "blog": f"Write a blog post about {topic}. Keep it concise but informative.",
        "social": f"Create 3 social media posts about {topic} for Twitter, LinkedIn, and Instagram.",
        "email": f"Write a short marketing email about {topic} with subject line and CTA.",
        "ppt": f"Create a PowerPoint presentation on '{topic}' with 6 slides. Each slide should have a title and 3–5 bullet points. Format:\nSlide 1: Title\n- Bullet\n- Bullet"
    }

    tones = {
        "professional": "Use a professional and authoritative tone.",
        "friendly": "Use a warm, conversational tone.",
        "persuasive": "Use a persuasive tone that creates urgency."
    }

    prompt = f"{prompts.get(content_type, prompts['blog'])}\n{tones.get(tone, tones['professional'])}"

    messages = [
        {"role": "system", "content": "You are an expert content creator."},
        {"role": "user", "content": prompt}
    ]

    payload = {
        "model": "mistral-medium-latest",
        "messages": messages,
        "temperature": 0.7,
        "stream": True
    }

    try:
        response = requests.post(endpoint, json=payload, headers=HEADERS, stream=True)
        response.raise_for_status()

        print(f"\nGenerating {tone} {content_type} content about '{topic}'...\n--- CONTENT ---\n")
        full_response = []

        for chunk in response.iter_lines():
            if chunk:
                line = chunk.decode('utf-8')
                if line.startswith('data: '):
                    line = line[6:]
                if line == '[DONE]':
                    break
                try:
                    content = json.loads(line)['choices'][0]['delta'].get('content', '')
                    if content:
                        print(content, end='', flush=True)
                        full_response.append(content)
                except json.JSONDecodeError:
                    continue

        print("\n")
        return ''.join(full_response)

    except Exception as e:
        print(f"❌ Error: {e}")
        return None

def save_to_pptx(slide_text, filename):
    """Convert AI output into a .pptx presentation file"""
    prs = Presentation()
    slide_chunks = [s.strip() for s in slide_text.split("Slide") if s.strip()]

    for chunk in slide_chunks:
        lines = chunk.split("\n")
        if not lines:
            continue
        title = lines[0].replace(":", "").strip()
        bullets = [line.strip("-• ").strip() for line in lines[1:] if line.strip()]

        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        content = slide.placeholders[1]

        for i, bullet in enumerate(bullets):
            if i == 0:
                content.text = bullet
            else:
                p = content.text_frame.add_paragraph()
                p.text = bullet

    pptx_filename = f"{filename}.pptx"

    prs.save(pptx_filename)
    print(f"📊 PPT saved to: {pptx_filename}")
    return pptx_filename  # Add this line!


def main():
    parser = argparse.ArgumentParser(description="Generate AI content with Mistral")
    parser.add_argument("topic", help="Topic for the content")
    parser.add_argument("--type", choices=["blog", "social", "email", "ppt"], default="blog", help="Content type")
    parser.add_argument("--tone", choices=["professional", "friendly", "persuasive"], default="professional", help="Content tone")

    args = parser.parse_args()

    content = generate_content_streaming(args.topic, args.type, args.tone)

    if content:
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        base_filename = f"{args.type}_{args.topic.replace(' ', '_')}_{timestamp}"

        # Save raw text
        with open(f"{base_filename}.txt", "w", encoding="utf-8") as f:
            f.write(content)
        print(f"📝 Text saved to: {base_filename}.txt")

        # Optional: Save .pptx
        if args.type == "ppt":
            save_to_pptx(content, base_filename)

if __name__ == "__main__":
    main()
